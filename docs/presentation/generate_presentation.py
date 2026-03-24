from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "docs" / "presentation" / "quant_research_platform_presentation.pptx"
SCREENS = REPO / "docs" / "screenshots"

BG = RGBColor(11, 15, 25)
SURFACE = RGBColor(17, 24, 39)
SURFACE_LIGHT = RGBColor(31, 41, 55)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(148, 163, 184)
BLUE = RGBColor(56, 189, 248)
GREEN = RGBColor(163, 230, 53)
PURPLE = RGBColor(192, 132, 252)
RED = RGBColor(248, 113, 113)


def base_slide(prs, slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.14))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    return slide


def add_title(slide, title, accent=None, subtitle=None):
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12.1), Inches(1.0))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE
    if accent:
        r2 = p.add_run()
        r2.text = accent
        r2.font.name = "Aptos Display"
        r2.font.size = Pt(28)
        r2.font.bold = True
        r2.font.color.rgb = GREEN
    if subtitle:
        sp = tf.add_paragraph()
        sp.text = subtitle
        sp.font.name = "Aptos"
        sp.font.size = Pt(12)
        sp.font.color.rgb = MUTED


def add_bullets(slide, left, top, width, height, items, color=MUTED, level0_color=WHITE, size=20):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.text = text
        p.level = level
        p.font.name = "Aptos"
        p.font.size = Pt(size if level == 0 else size - 3)
        p.font.color.rgb = level0_color if level == 0 else color
        p.space_after = Pt(8)
    return box


def add_panel(slide, left, top, width, height, title=None, border=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    shape.line.color.rgb = border
    shape.line.width = Pt(1.2)
    shape.adjustments[0] = 0.08
    if title:
        t = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.3))
        p = t.text_frame.paragraphs[0]
        p.text = title
        p.font.name = "Aptos Display"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
    return shape


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = base_slide(prs, prs.slides.add_slide(prs.slide_layouts[6]))
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(7.3), Inches(3.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Quantitative Research Platform"
    r.font.name = "Aptos Display"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "for European S&P 500 Options"
    r2.font.name = "Aptos Display"
    r2.font.size = Pt(30)
    r2.font.bold = True
    r2.font.color.rgb = GREEN
    p3 = tf.add_paragraph()
    p3.text = "Combining pricing models, calibration, live valuation-gap scanning, and controlled backtesting."
    p3.font.name = "Aptos"
    p3.font.size = Pt(18)
    p3.font.color.rgb = MUTED
    p3.space_before = Pt(18)
    p4 = tf.add_paragraph()
    p4.text = "Bach Nguyen"
    p4.font.name = "Aptos"
    p4.font.size = Pt(18)
    p4.font.bold = True
    p4.font.color.rgb = WHITE
    p4.space_before = Pt(24)
    p5 = tf.add_paragraph()
    p5.text = "github.com/bachnguyennn/black-scholes-monte-carlo"
    p5.font.name = "Aptos"
    p5.font.size = Pt(13)
    p5.font.color.rgb = BLUE
    img = SCREENS / "option-pricing.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(8.3), Inches(1.0), width=Inches(4.2), height=Inches(4.8))

    slide = base_slide(prs, prs.slides.add_slide(prs.slide_layouts[6]))
    add_title(slide, "Problem & ", "Motivation")
    add_panel(slide, Inches(0.6), Inches(1.3), Inches(7.5), Inches(5.5), border=BLUE)
    add_bullets(slide, Inches(0.9), Inches(1.7), Inches(6.9), Inches(4.8), [
        "European index options are central to institutional hedging and volatility trading.",
        "Black-Scholes is useful, but its assumptions are too restrictive for real markets.",
        "Real option markets exhibit skew, smile, stochastic volatility, and jump risk.",
        "The objective was to build a platform that compares models and makes assumptions explicit.",
    ], size=20)
    add_panel(slide, Inches(8.5), Inches(1.7), Inches(4.1), Inches(4.2), title="Why This Matters", border=GREEN)
    add_bullets(slide, Inches(8.8), Inches(2.25), Inches(3.5), Inches(3.2), [
        "Better pricing intuition",
        "Transparent model comparison",
        "Calibration tied to live data",
        "Research workflow instead of black-box output",
    ], level0_color=WHITE, size=18)

    slide = base_slide(prs, prs.slides.add_slide(prs.slide_layouts[6]))
    add_title(slide, "System ", "Architecture")
    for x, title, subtitle, color in [
        (0.7, "Frontend", "Streamlit\nInteraction, charts, diagnostics", BLUE),
        (4.55, "Backend", "FastAPI\nAPI endpoints and scan execution", GREEN),
        (8.4, "Core Engines", "Pricing, calibration, scanner, backtester", PURPLE),
    ]:
        add_panel(slide, Inches(x), Inches(2.0), Inches(3.2), Inches(2.8), title=title, border=color)
        t = slide.shapes.add_textbox(Inches(x) + Inches(0.22), Inches(2.7), Inches(2.75), Inches(1.5))
        p = t.text_frame.paragraphs[0]
        for idx, line in enumerate(subtitle.split("\n")):
            if idx == 0:
                p.text = line
                p.font.name = "Aptos Display"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = color
            else:
                p2 = t.text_frame.add_paragraph()
                p2.text = line
                p2.font.name = "Aptos"
                p2.font.size = Pt(16)
                p2.font.color.rgb = MUTED
    for x in [3.82, 7.67]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(3.05), Inches(0.42), Inches(0.42))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = SURFACE_LIGHT
        arrow.line.fill.background()

    slide = base_slide(prs, prs.slides.add_slide(prs.slide_layouts[6]))
    add_title(slide, "Quantitative Models & ", "Math")
    positions = [
        (0.7, 1.6, 5.8, 1.0, "Black-Scholes", "Analytical benchmark for European options.", BLUE),
        (0.7, 2.85, 5.8, 1.0, "GBM Monte Carlo", "Baseline simulation engine and convergence reference.", BLUE),
        (0.7, 4.1, 5.8, 1.0, "Heston", "Captures stochastic volatility and volatility skew.", GREEN),
        (6.8, 1.6, 5.8, 1.0, "Jump Diffusion", "Adds discontinuous crash risk to diffusion dynamics.", RED),
        (6.8, 2.85, 5.8, 1.0, "LSV", "Extends the framework toward better surface fitting.", PURPLE),
    ]
    for x, y, w, h, title, desc, color in positions:
        add_panel(slide, Inches(x), Inches(y), Inches(w), Inches(h), border=color)
        tb = slide.shapes.add_textbox(Inches(x) + Inches(0.2), Inches(y) + Inches(0.2), Inches(w) - Inches(0.4), Inches(h) - Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.name = "Aptos Display"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p2 = tb.text_frame.add_paragraph()
        p2.text = desc
        p2.font.name = "Aptos"
        p2.font.size = Pt(14)
        p2.font.color.rgb = MUTED

    slide = base_slide(prs, prs.slides.add_slide(prs.slide_layouts[6]))
    add_title(slide, "Calibration: ", "Connecting Theory to Data")
    add_panel(slide, Inches(0.7), Inches(1.7), Inches(3.0), Inches(3.8), border=PURPLE)
    num = slide.shapes.add_textbox(Inches(1.05), Inches(2.2), Inches(2.3), Inches(1.3))
    p = num.text_frame.paragraphs[0]
    p.text = "RMSE"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos Display"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p2 = num.text_frame.add_paragraph()
    p2.text = "Root Mean Square Error"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Aptos"
    p2.font.size = Pt(16)
    p2.font.color.rgb = MUTED
    add_panel(slide, Inches(4.1), Inches(1.7), Inches(8.5), Inches(4.2), title="Why Calibration Matters", border=BLUE)
    add_bullets(slide, Inches(4.4), Inches(2.25), Inches(7.8), Inches(3.2), [
        "Without calibration, advanced model outputs are only parameter guesses.",
        "Calibration fits model parameters to the implied volatility surface.",
        "Weighted error minimization aligns theoretical values with live market structure.",
        "This is the bridge between pure theory and observed option prices.",
    ], size=18)

    slide = base_slide(prs, prs.slides.add_slide(prs.slide_layouts[6]))
    add_title(slide, "Live ", "Scanner Logic")
    add_panel(slide, Inches(0.7), Inches(1.5), Inches(5.4), Inches(4.9), border=BLUE)
    add_bullets(slide, Inches(1.0), Inches(1.95), Inches(4.8), Inches(4.1), [
        "Strict filtering removes illiquid, inconsistent, and low-quality contracts.",
        "Signal logic is bid/ask aware, not midpoint based.",
        "Opportunities are only shown when model edge exceeds execution costs.",
        "Diagnostics expose filtering reasons, volatility source, and execution mode.",
    ], size=18)
    img = SCREENS / "live-scanner.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(6.35), Inches(1.55), width=Inches(6.25), height=Inches(4.7))

    slide = base_slide(prs, prs.slides.add_slide(prs.slide_layouts[6]))
    add_title(slide, "Backtester & ", "Research Integrity")
    img = SCREENS / "backtester.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(0.75), Inches(1.65), width=Inches(6.1), height=Inches(4.7))
    add_panel(slide, Inches(7.1), Inches(1.65), Inches(5.5), Inches(4.8), border=GREEN)
    add_bullets(slide, Inches(7.4), Inches(2.0), Inches(4.8), Inches(4.1), [
        "No-look-ahead volatility estimation.",
        "Delta hedging and transaction-cost accounting.",
        "Cost sensitivity and methodology disclosure.",
        "Research simulator framing rather than execution-grade claim.",
    ], size=18)

    slide = base_slide(prs, prs.slides.add_slide(prs.slide_layouts[6]))
    add_title(slide, "Results & ", "Future Work")
    add_panel(slide, Inches(0.7), Inches(1.55), Inches(3.8), Inches(4.95), title="Strengths", border=GREEN)
    add_bullets(slide, Inches(0.95), Inches(2.05), Inches(3.1), Inches(3.9), [
        "Multi-model pricing framework",
        "Calibration layer",
        "Transparent diagnostics",
        "Modular architecture",
    ], size=17)
    add_panel(slide, Inches(4.8), Inches(1.55), Inches(3.6), Inches(4.95), title="Limitations", border=RED)
    add_bullets(slide, Inches(5.05), Inches(2.05), Inches(2.95), Inches(3.9), [
        "Depends on Yahoo Finance data",
        "No historical option NBBO replay",
        "Research-grade, not execution-grade",
    ], size=17)
    add_panel(slide, Inches(8.7), Inches(1.55), Inches(3.9), Inches(4.95), title="Future Work", border=BLUE)
    add_bullets(slide, Inches(8.95), Inches(2.05), Inches(3.2), Inches(3.9), [
        "Institutional data feed",
        "Historical options dataset",
        "Stronger calibration validation",
        "Automated reporting",
    ], size=17)
    footer = slide.shapes.add_textbox(Inches(0.8), Inches(6.7), Inches(11.8), Inches(0.35))
    p = footer.text_frame.paragraphs[0]
    p.text = "A credible quantitative research environment combining model comparison, calibration, diagnostics, and transparent backtesting."
    p.font.name = "Aptos"
    p.font.size = Pt(14)
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build_deck()
